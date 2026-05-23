from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt


ROOT = Path(__file__).resolve().parent
REPORT_MD = ROOT / "RAPPORT_STAGE_EDUPRIMAIRE.md"
DOCX_OUT = ROOT / "RAPPORT_STAGE_EDUPRIMAIRE.docx"
PPTX_OUT = ROOT / "PRESENTATION_EDUPRIMAIRE_7_SLIDES.pptx"


def add_markdown_to_docx(md_path: Path, out_path: Path) -> None:
    doc = Document()
    styles = doc.styles
    styles["Normal"].font.name = "Calibri"
    styles["Normal"].font.size = Pt(11)

    in_code = False
    for raw in md_path.read_text(encoding="utf-8").splitlines():
        line = raw.rstrip()

        if line.startswith("```"):
            in_code = not in_code
            if in_code:
                p = doc.add_paragraph()
                r = p.add_run("Diagramme / bloc technique :")
                r.bold = True
            continue

        if not line:
            continue

        if in_code:
            p = doc.add_paragraph()
            r = p.add_run(line)
            r.font.name = "Courier New"
            r.font.size = Pt(8)
            continue

        if line.startswith("# "):
            title = doc.add_heading(line[2:], level=0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            continue
        if line.startswith("## "):
            doc.add_heading(line[3:], level=1)
            continue
        if line.startswith("### "):
            doc.add_heading(line[4:], level=2)
            continue
        if line.startswith("#### "):
            doc.add_heading(line[5:], level=3)
            continue
        if line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
            continue
        if line.startswith("|"):
            p = doc.add_paragraph()
            r = p.add_run(line)
            r.font.name = "Courier New"
            r.font.size = Pt(8)
            continue
        if line.startswith("---"):
            continue

        p = doc.add_paragraph()
        if line.startswith("**") and line.endswith("**"):
            run = p.add_run(line.strip("*"))
            run.bold = True
        else:
            p.add_run(line.replace("**", ""))

    section = doc.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)
    doc.save(out_path)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptRGBColor(*color)


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(8.9), Inches(0.75))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptPt(30)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(30, 58, 95)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.68), Inches(1.25), Inches(8.6), Inches(0.35))
        stf = sub.text_frame
        stf.text = subtitle
        stf.paragraphs[0].font.size = PptPt(13)
        stf.paragraphs[0].font.color.rgb = PptRGBColor(87, 100, 116)


def add_bullets(slide, bullets, x=0.8, y=1.75, w=8.5, h=4.8, font_size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = PptPt(font_size)
        p.font.color.rgb = PptRGBColor(34, 40, 49)
        p.space_after = PptPt(8)


def add_footer(slide, index):
    box = slide.shapes.add_textbox(Inches(8.6), Inches(6.95), Inches(1.0), Inches(0.25))
    tf = box.text_frame
    tf.text = f"{index}/7"
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf.paragraphs[0].font.size = PptPt(9)
    tf.paragraphs[0].font.color.rgb = PptRGBColor(120, 132, 148)


def add_architecture(slide):
    labels = ["Navigateur", "React + Vite", "API Django REST", "SQLite"]
    xs = [0.8, 3.0, 5.35, 7.75]
    for x, label in zip(xs, labels):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(2.7), Inches(1.55), Inches(0.8))
        set_fill(shape, (30, 58, 95))
        shape.line.color.rgb = PptRGBColor(30, 58, 95)
        tf = shape.text_frame
        tf.text = label
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = PptPt(13)
        tf.paragraphs[0].font.color.rgb = PptRGBColor(255, 255, 255)

    for x in [2.4, 4.75, 7.15]:
        line = slide.shapes.add_connector(1, Inches(x), Inches(3.1), Inches(x + 0.45), Inches(3.1))
        line.line.color.rgb = PptRGBColor(243, 183, 73)
        line.line.width = PptPt(3)


def make_pptx(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = [
        (
            "EduPrimaire",
            "Application web de gestion d'école primaire",
            [
                "Projet de soutenance - Licence en informatique",
                "Gestion administrative, pédagogique et financière",
                "Frontend React, Backend Django REST, Base SQLite",
            ],
        ),
        (
            "Contexte et problématique",
            None,
            [
                "Gestion encore souvent manuelle : cahiers, fichiers séparés, traitements répétitifs.",
                "Risque d'erreurs, doublons, pertes d'informations et lenteur des traitements.",
                "Difficulté à suivre les présences, paiements, bulletins et échanges avec les parents.",
                "Besoin : centraliser et sécuriser la gestion d'une école primaire.",
            ],
        ),
        (
            "Objectifs du projet",
            None,
            [
                "Gérer élèves, parents, enseignants, classes et matières.",
                "Suivre notes, bulletins, présences et emplois du temps.",
                "Gérer inscriptions, frais de scolarité et paiements.",
                "Proposer messagerie, notifications et portail parent.",
                "Sécuriser l'accès selon les rôles utilisateurs.",
            ],
        ),
        (
            "Architecture et technologies",
            "Architecture client-serveur avec API REST",
            [
                "Frontend : React 18, Vite, Tailwind CSS, React Router.",
                "Backend : Python, Django 4.2, Django REST Framework.",
                "Sécurité : authentification JWT avec refresh token.",
                "Production : Gunicorn, WhiteNoise, Render.",
            ],
        ),
        (
            "Modélisation UML",
            None,
            [
                "Acteurs : administrateur, directeur, enseignant, secrétaire, comptable, parent.",
                "Entités : école, utilisateur, classe, matière, élève, parent, enseignant.",
                "Suivi : notes, bulletins, présences, scolarité, paiements, messages, notifications.",
                "Diagrammes réalisés : cas d'utilisation, domaine, classes, packages, séquences.",
            ],
        ),
        (
            "Fonctionnalités réalisées",
            None,
            [
                "Tableau de bord et navigation adaptée au rôle.",
                "Gestion des élèves, parents, enseignants, classes et matières.",
                "Notes, bulletins, présences et emploi du temps.",
                "Secrétariat, inscriptions, scolarité et paiements.",
                "Messagerie, bibliothèque, notifications et portail parent.",
            ],
        ),
        (
            "Bilan et perspectives",
            None,
            [
                "Application fonctionnelle, modulaire et prête pour déploiement backend.",
                "Séparation claire entre frontend React et backend Django REST.",
                "Perspectives : PostgreSQL, tests automatisés, statistiques avancées.",
                "Évolutions possibles : email/SMS, PWA complète, portail parent enrichi.",
            ],
        ),
    ]

    for index, (title, subtitle, bullets) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = PptRGBColor(248, 247, 242)

        accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
        set_fill(accent, (243, 183, 73))
        accent.line.fill.background()

        add_title(slide, title, subtitle)
        if index == 4:
            add_bullets(slide, bullets, y=1.75, h=1.8, font_size=15)
            add_architecture(slide)
        else:
            add_bullets(slide, bullets, font_size=18 if index != 1 else 19)
        add_footer(slide, index)

    prs.save(out_path)


if __name__ == "__main__":
    add_markdown_to_docx(REPORT_MD, DOCX_OUT)
    make_pptx(PPTX_OUT)
    print(DOCX_OUT)
    print(PPTX_OUT)
